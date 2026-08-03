#!/usr/bin/env python3
"""Build a polished, idea-led PowerPoint from the living survey (English)."""

from __future__ import annotations

from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = Path(__file__).resolve().parent / "out"
PPTX_NAME = "next-gen-ai-compiler-survey.pptx"

# Editorial palette: deep ink, paper, single copper accent (no purple/cream clichés).
INK = RGBColor(0x12, 0x16, 0x1C)
PAPER = RGBColor(0xF7, 0xF5, 0xF1)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xD9, 0xD4, 0xCC)
MUTED = RGBColor(0x6A, 0x66, 0x5E)
BODY = RGBColor(0x2A, 0x28, 0x24)
ACCENT = RGBColor(0xB5, 0x4A, 0x24)  # copper
TEAL = RGBColor(0x1F, 0x5C, 0x56)
SOFT_TEAL = RGBColor(0xE4, 0xEF, 0xED)
SOFT_COPPER = RGBColor(0xF6, 0xEB, 0xE4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

DISPLAY = "Noto Serif"
SANS = "Inter"
W, H = 13.333, 7.5


def font(run, size=16, bold=False, color=BODY, name=SANS):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name


def blank(prs: Presentation):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(W), Inches(H))
    bg.fill.solid()
    bg.fill.fore_color.rgb = PAPER
    bg.line.fill.background()
    return s


def rect(slide, l, t, w, h, fill, line=None):
    sh = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h)
    )
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
    return sh


def round_rect(slide, l, t, w, h, fill, line=None):
    sh = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h)
    )
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
    try:
        sh.adjustments[0] = 0.08
    except Exception:
        pass
    return sh


def textbox(slide, l, t, w, h, text, size=16, bold=False, color=BODY, name=SANS, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    font(r, size, bold, color, name)
    return box


def multiline(slide, l, t, w, h, lines, size=14, color=BODY, name=SANS, bold_first=False, gap=6):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        r = p.add_run()
        r.text = line
        font(r, size, bold_first and i == 0, color, name)
    return box


def kicker(slide, l, t, label):
    textbox(slide, l, t, 4, 0.3, label.upper(), 11, True, ACCENT, SANS)


def page_footer(slide, n, total):
    textbox(slide, 0.6, 7.1, 8, 0.25, "Next-Gen AI Compiler Survey", 10, False, MUTED, SANS)
    textbox(slide, 11.2, 7.1, 1.5, 0.25, f"{n} / {total}", 10, False, MUTED, SANS, PP_ALIGN.RIGHT)


def add_chart(slide, chart_type, l, t, w, h, cats, series, legend=False):
    data = CategoryChartData()
    data.categories = cats
    for name, vals in series.items():
        data.add_series(name, vals)
    chart = slide.shapes.add_chart(
        chart_type, Inches(l), Inches(t), Inches(w), Inches(h), data
    ).chart
    chart.has_legend = legend
    if legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
    colors = [TEAL, ACCENT, MUTED, INK]
    try:
        for i, ser in enumerate(chart.series):
            ser.format.fill.solid()
            ser.format.fill.fore_color.rgb = colors[i % len(colors)]
    except Exception:
        pass
    return chart


# ── slides ──────────────────────────────────────────────────────────

def s_title(prs):
    s = blank(prs)
    rect(s, 0, 0, W, H, INK)
    rect(s, 0, 0, 0.18, H, ACCENT)
    textbox(s, 0.9, 1.8, 11, 0.4, "LIVING SURVEY", 12, True, ACCENT, SANS)
    textbox(
        s, 0.9, 2.3, 11.5, 1.4,
        "The next compiler\nis agentic — not replaced.",
        40, True, WHITE, DISPLAY,
    )
    textbox(
        s, 0.9, 4.0, 10.5, 1.0,
        "A prediction for ~2027–28 and the next five years:\nhow agents reshape the AI compiler stack and HW–SW codesign.",
        18, False, RGBColor(0xC8, 0xC3, 0xBA), SANS,
    )
    textbox(
        s, 0.9, 6.4, 10, 0.4,
        f"Generated {date.today().isoformat()}  ·  From docs/SURVEY · ROADMAP · STACK · CLAIMS",
        12, False, MUTED, SANS,
    )


def s_thesis(prs, n, total):
    s = blank(prs)
    kicker(s, 0.7, 0.45, "01  —  The thesis")
    textbox(s, 0.7, 0.95, 12, 1.0, "Agents search. Compilers decide.", 32, True, INK, DISPLAY)
    textbox(
        s, 0.7, 2.0, 11.5, 0.8,
        "The winning pattern across ACCLAIM, HintPilot, AgentCompile, Magellan, TritorX, and GEAK is hybrid:",
        16, False, MUTED, SANS,
    )
    # two big quote cards
    left = round_rect(s, 0.7, 3.1, 5.7, 3.0, SOFT_TEAL)
    right = round_rect(s, 6.9, 3.1, 5.7, 3.0, SOFT_COPPER)
    textbox(s, 1.0, 3.4, 5.1, 0.4, "Agents own", 14, True, TEAL, SANS)
    multiline(
        s, 1.0, 3.9, 5.1, 1.8,
        ["semantic search", "orchestration & budgets", "artifact synthesis", "bring-up coverage loops"],
        18, INK, DISPLAY,
    )
    textbox(s, 7.2, 3.4, 5.1, 0.4, "Compilers own", 14, True, ACCENT, SANS)
    multiline(
        s, 7.2, 3.9, 5.1, 1.8,
        ["lowering & legality", "measurement oracles", "admit / fallback", "deterministic defaults"],
        18, INK, DISPLAY,
    )
    page_footer(s, n, total)


def s_not_replace(prs, n, total):
    s = blank(prs)
    kicker(s, 0.7, 0.45, "02  —  Hard limit")
    textbox(s, 0.7, 0.95, 12, 1.2, "Free-form IR rewrite\nkeeps losing.", 32, True, INK, DISPLAY)
    # evidence strip
    cards = [
        ("mlirAgent", "Frontier models scored\nbelow identity on IR transforms"),
        ("HintPilot / AgentCompile", "Succeed by constraining\nactions to hints & templates"),
        ("ACCLAIM", "Best when compiler tools\nand tests remain admit gates"),
    ]
    for i, (t, b) in enumerate(cards):
        x = 0.7 + i * 4.1
        round_rect(s, x, 3.5, 3.9, 2.6, CARD, LINE)
        rect(s, x, 3.5, 3.9, 0.08, ACCENT if i == 0 else TEAL)
        textbox(s, x + 0.3, 3.8, 3.3, 0.5, t, 16, True, INK, DISPLAY)
        multiline(s, x + 0.3, 4.5, 3.3, 1.4, b.split("\n"), 14, MUTED, SANS)
    page_footer(s, n, total)


def s_architecture(prs, n, total):
    s = blank(prs)
    kicker(s, 0.7, 0.4, "03  —  Architecture")
    textbox(s, 0.7, 0.85, 12, 0.6, "Three planes, one contract.", 28, True, INK, DISPLAY)

    planes = [
        (1.3, TEAL, "Control plane", "Propose · measure · admit\nFour agent jobs (a–d)"),
        (3.2, INK, "Data plane", "Lower · verify · fallback\nInductor / MLIR / Triton / Tile"),
        (5.1, ACCENT, "Codesign feedback", "Sim + silicon traces\n→ ISA / dialect RFCs"),
    ]
    for top, color, title, body in planes:
        round_rect(s, 1.5, top, 10.3, 1.55, CARD, LINE)
        rect(s, 1.5, top, 0.16, 1.55, color)
        textbox(s, 2.0, top + 0.25, 4, 0.4, title, 18, True, INK, DISPLAY)
        multiline(s, 6.5, top + 0.3, 5, 1.1, body.split("\n"), 15, MUTED, SANS)
    page_footer(s, n, total)


def s_four_jobs(prs, n, total):
    s = blank(prs)
    kicker(s, 0.7, 0.35, "04  —  Four jobs")
    textbox(s, 0.7, 0.75, 12, 0.55, "What agents actually do in compilers.", 28, True, INK, DISPLAY)
    jobs = [
        ("a", "Online specialize", "At compile / serve time", "Hints, ACFs, Triton kernels,\npass lists for this workload", "CompileIQ · GEAK · AutoKernel · ACCLAIM"),
        ("b", "Offline evolve", "Compiler engineering time", "Readable C++ heuristics\nand MLGO features", "Magellan · AlphaEvolve · OpenEvolve"),
        ("c", "Oracle review", "PR / change review", "Alive2 / opt-gated agents\nbeat generic chat review", "Archer · LLVM agent review"),
        ("d", "Bring-up / codesign", "New ASIC / NPU TTM", "Coverage then perf on\nsim + silicon", "TritorX · KernelEvolve · Ascend diag."),
    ]
    for i, (letter, title, when, what, who) in enumerate(jobs):
        x = 0.45 + (i % 2) * 6.4
        y = 1.55 + (i // 2) * 2.55
        round_rect(s, x, y, 6.15, 2.35, CARD, LINE)
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.25), Inches(y + 0.3), Inches(0.55), Inches(0.55))
        circle.fill.solid()
        circle.fill.fore_color.rgb = ACCENT if i % 2 else TEAL
        circle.line.fill.background()
        tf = circle.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        r = tf.paragraphs[0].add_run()
        r.text = letter
        font(r, 16, True, WHITE, SANS)
        textbox(s, x + 1.0, y + 0.28, 4.8, 0.4, title, 18, True, INK, DISPLAY)
        textbox(s, x + 1.0, y + 0.7, 4.8, 0.3, when, 12, False, ACCENT, SANS)
        multiline(s, x + 0.35, y + 1.15, 5.5, 0.7, what.split("\n"), 13, BODY, SANS)
        textbox(s, x + 0.35, y + 1.9, 5.5, 0.3, who, 11, False, MUTED, SANS)
    page_footer(s, n, total)


def s_era(prs, n, total):
    s = blank(prs)
    kicker(s, 0.7, 0.4, "05  —  How we got here")
    textbox(s, 0.7, 0.85, 12, 0.55, "From autotune gyms to agent loops.", 28, True, INK, DISPLAY)
    eras = [
        ("2018–22", "DL compilers mature", "TVM, Ansor, XLA, MLIR"),
        ("2020–23", "RL for compilers", "CompilerGym, MLGO advisors"),
        ("2023–24", "LLMs enter IR", "Pass lists, Meta LLM Compiler"),
        ("2025–26", "Agentic hybrid", "Kernels, heuristics, bring-up"),
    ]
    for i, (y, title, body) in enumerate(eras):
        x = 0.7 + i * 3.15
        # timeline node
        node = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 1.2), Inches(2.0), Inches(0.28), Inches(0.28))
        node.fill.solid()
        node.fill.fore_color.rgb = ACCENT if i == 3 else TEAL
        node.line.fill.background()
        if i < 3:
            rect(s, x + 1.48, 2.1, 2.9, 0.04, LINE)
        textbox(s, x, 2.5, 3.0, 0.35, y, 12, True, ACCENT, SANS)
        textbox(s, x, 2.95, 3.0, 0.7, title, 16, True, INK, DISPLAY)
        textbox(s, x, 3.7, 3.0, 0.8, body, 13, False, MUTED, SANS)
    # insight
    round_rect(s, 0.7, 5.0, 12.0, 1.5, SOFT_TEAL)
    textbox(
        s, 1.0, 5.35, 11.4, 1.0,
        "Insight: each era kept the data plane and moved intelligence into a better control interface —\ncost models → RL policies → foundation LLMs → tool-using agents with oracles.",
        15, False, INK, SANS,
    )
    page_footer(s, n, total)


def s_ideas_grid(prs, n, total):
    s = blank(prs)
    kicker(s, 0.7, 0.35, "06  —  Ideas the survey surfaces")
    textbox(s, 0.7, 0.75, 12, 0.5, "Six ideas worth stealing.", 28, True, INK, DISPLAY)
    ideas = [
        ("Constrained action spaces", "Hints, templates, EVOLVE-blocks, ACF knobs — not free codegen."),
        ("Admit gates as product", "Tests, Alive2, OpInfo, profilers are the compiler’s API to agents."),
        ("Artifacts over binaries", "ACFs, heuristics, memories, bring-up corpora belong in VCS."),
        ("Coverage → perf ladder", "TritorX then KernelEvolve: new silicon needs ops before peak."),
        ("Negative results are Tier A", "mlirAgent’s below-identity bound shapes the architecture."),
        ("Conflicts are load-bearing", "Magellan vs MLGO, vendor vs KernelBench-X — keep both sides."),
    ]
    for i, (t, b) in enumerate(ideas):
        x = 0.55 + (i % 3) * 4.2
        y = 1.55 + (i // 3) * 2.5
        round_rect(s, x, y, 4.0, 2.25, CARD, LINE)
        textbox(s, x + 0.25, y + 0.35, 3.5, 0.7, t, 16, True, INK, DISPLAY)
        multiline(s, x + 0.25, y + 1.15, 3.5, 0.9, [b], 13, MUTED, SANS)
    page_footer(s, n, total)


def s_stack(prs, n, total):
    s = blank(prs)
    kicker(s, 0.7, 0.4, "07  —  Stack reshape")
    textbox(s, 0.7, 0.85, 12, 0.55, "Where the agentic compiler presses.", 28, True, INK, DISPLAY)
    layers = [
        ("Framework", 7, "Amdahl-rank hot regions; write kernels back"),
        ("Kernel DSL", 9, "Triton / Helion / Tile become agent surfaces"),
        ("Portable IR", 6, "Fingerprints + tools; free rewrite fails"),
        ("Mid / back end", 8, "Heuristics, passes, ACFs"),
        ("Oracles", 9, "Admit gates & rewards"),
        ("Artifacts", 8, "First-class in CI / review"),
        ("Serving", 7, "Specialize without breaking graphs"),
        ("Silicon / sim", 8, "Bring-up + ISA feedback"),
    ]
    # horizontal bars
    for i, (name, score, note) in enumerate(layers):
        y = 1.55 + i * 0.62
        textbox(s, 0.7, y, 2.2, 0.4, name, 13, True, INK, SANS)
        track = rect(s, 3.0, y + 0.08, 6.5, 0.28, LINE)
        bar = rect(s, 3.0, y + 0.08, 6.5 * (score / 10), 0.28, TEAL if score < 9 else ACCENT)
        textbox(s, 9.7, y, 3.2, 0.45, note, 11, False, MUTED, SANS)
    page_footer(s, n, total)


def s_conflicts(prs, n, total):
    s = blank(prs)
    kicker(s, 0.7, 0.35, "08  —  Productive disagreements")
    textbox(s, 0.7, 0.75, 12, 0.55, "Do not average these away.", 28, True, INK, DISPLAY)
    rows = [
        ("C1", "Magellan heuristics", "MLGO neural advisors", "Parallel production bets"),
        ("C2", "Vendor speedup headlines", "KernelBench-X ceilings", "Need p50/p90 traces"),
        ("C3", "Wide rewrite APIs", "Narrow advisory APIs", "Oracles decide"),
        ("C9", "Coverage-first bring-up", "Peak-perf kernel agents", "Ladder, not either/or"),
        ("C10", "Autonomous chip design", "Compiler codesign feedback", "We bet on feedback only"),
    ]
    # header
    textbox(s, 0.7, 1.5, 1.2, 0.35, "ID", 11, True, MUTED, SANS)
    textbox(s, 1.9, 1.5, 3.5, 0.35, "Claim A", 11, True, MUTED, SANS)
    textbox(s, 5.6, 1.5, 3.5, 0.35, "Claim B", 11, True, MUTED, SANS)
    textbox(s, 9.3, 1.5, 3.5, 0.35, "Stance", 11, True, MUTED, SANS)
    rect(s, 0.7, 1.9, 12.0, 0.02, LINE)
    for i, (cid, a, b, stance) in enumerate(rows):
        y = 2.1 + i * 0.85
        round_rect(s, 0.7, y, 12.0, 0.75, CARD if i % 2 == 0 else SOFT_TEAL, LINE)
        textbox(s, 0.9, y + 0.2, 1.0, 0.4, cid, 14, True, ACCENT, SANS)
        textbox(s, 1.9, y + 0.2, 3.5, 0.4, a, 13, False, INK, SANS)
        textbox(s, 5.6, y + 0.2, 3.5, 0.4, b, 13, False, INK, SANS)
        textbox(s, 9.3, y + 0.2, 3.2, 0.4, stance, 13, True, TEAL, SANS)
    page_footer(s, n, total)


def s_codesign(prs, n, total):
    s = blank(prs)
    kicker(s, 0.7, 0.4, "09  —  HW–SW codesign")
    textbox(s, 0.7, 0.85, 12, 0.7, "Still an agentic-compiler problem.", 28, True, INK, DISPLAY)
    textbox(
        s, 0.7, 1.6, 12, 0.5,
        "Custom ASICs win or lose on operator coverage × generations × devices. Agents eat that matrix.",
        15, False, MUTED, SANS,
    )
    steps = [
        ("Spec / sim", "Draft ISA +\nfuture-device sim"),
        ("Coverage", "TritorX-class\nATen / OpInfo"),
        ("Perf", "KernelEvolve /\nprofile search"),
        ("Serve", "e2e latency\n& TCO gates"),
        ("Feedback", "Traces → next\nISA / dialect"),
    ]
    for i, (t, b) in enumerate(steps):
        x = 0.55 + i * 2.55
        round_rect(s, x, 2.5, 2.35, 2.5, CARD, LINE)
        rect(s, x, 2.5, 2.35, 0.1, ACCENT if i == 4 else TEAL)
        textbox(s, x + 0.15, 2.85, 2.05, 0.7, f"{i+1}. {t}", 15, True, INK, DISPLAY)
        multiline(s, x + 0.15, 3.7, 2.05, 1.0, b.split("\n"), 13, MUTED, SANS)
        if i < 4:
            textbox(s, x + 2.15, 3.4, 0.4, 0.4, "→", 18, True, MUTED, SANS)
    round_rect(s, 0.7, 5.4, 12.0, 1.2, SOFT_COPPER)
    textbox(
        s, 1.0, 5.7, 11.4, 0.7,
        "Non-goal: autonomous tape-out. Agents stress compilers and file ISA pain; humans + EDA own silicon.",
        15, False, INK, SANS,
    )
    page_footer(s, n, total)


def s_roadmap(prs, n, total):
    s = blank(prs)
    kicker(s, 0.7, 0.35, "10  —  Roadmap")
    textbox(s, 0.7, 0.75, 12, 0.5, "What ships — and what does not.", 28, True, INK, DISPLAY)
    # two columns
    round_rect(s, 0.55, 1.5, 6.0, 5.1, CARD, LINE)
    rect(s, 0.55, 1.5, 6.0, 0.55, TEAL)
    textbox(s, 0.8, 1.6, 5.5, 0.4, "2027–28  ·  likely", 16, True, WHITE, SANS)
    multiline(
        s, 0.85, 2.3, 5.4, 4.0,
        [
            "Agent-addressable tool APIs in major stacks",
            "CI-gated specialize on hot kernels / apps",
            "Magellan-class and MLGO both still live",
            "Oracle PR review in serious compiler orgs",
            "Coverage→perf bring-up on new ASICs",
            "Triton-family still primary agent surface",
        ],
        14, BODY, SANS, gap=10,
    )
    round_rect(s, 6.85, 1.5, 6.0, 5.1, CARD, LINE)
    rect(s, 6.85, 1.5, 6.0, 0.55, ACCENT)
    textbox(s, 7.1, 1.6, 5.5, 0.4, "Not soon", 16, True, WHITE, SANS)
    multiline(
        s, 7.15, 2.3, 5.4, 4.0,
        [
            "LLM replaces opt / Inductor end-to-end",
            "One universal agent IR for all vendors",
            "Uniform wins on fusion-heavy public ladders",
            "Agents autonomously design microarchitecture",
            "Generic SCM chat as compiler evidence",
            "Silent default agents without distributions",
        ],
        14, BODY, SANS, gap=10,
    )
    page_footer(s, n, total)


def s_evidence_chart(prs, n, total):
    s = blank(prs)
    kicker(s, 0.7, 0.4, "11  —  Evidence shape")
    textbox(s, 0.7, 0.85, 12, 0.5, "A living bibliography, not a catalog.", 28, True, INK, DISPLAY)
    add_chart(
        s, XL_CHART_TYPE.BAR_CLUSTERED,
        0.5, 1.6, 7.5, 5.0,
        [
            "GPU kernels", "Agentic / RL", "SCM / review", "Company infra",
            "Classic DL", "Foundation LLMs", "Forums", "HW codesign",
        ],
        {"Digests": [16, 13, 13, 10, 9, 8, 8, 4]},
        legend=False,
    )
    round_rect(s, 8.3, 1.8, 4.5, 4.5, CARD, LINE)
    textbox(s, 8.6, 2.1, 4.0, 0.4, "How to read it", 16, True, INK, DISPLAY)
    multiline(
        s, 8.6, 2.7, 4.0, 3.2,
        [
            "★ digests drive ROADMAP",
            "Tier A reshapes compile",
            "Tier B is substrate only",
            "Tier C is demoted",
            "Conflicts beat false consensus",
            "Mechanisms > headline ×",
        ],
        14, MUTED, SANS, gap=8,
    )
    page_footer(s, n, total)


def s_claims(prs, n, total):
    s = blank(prs)
    kicker(s, 0.7, 0.4, "12  —  Claim confidence")
    textbox(s, 0.7, 0.85, 12, 0.5, "What we believe — for now.", 28, True, INK, DISPLAY)
    add_chart(
        s, XL_CHART_TYPE.COLUMN_CLUSTERED,
        0.6, 1.6, 7.8, 5.0,
        ["Supported", "Contested", "Watch"],
        {"Claims": [10, 2, 2]},
        legend=False,
    )
    round_rect(s, 8.7, 1.8, 4.1, 4.5, SOFT_COPPER)
    textbox(s, 9.0, 2.2, 3.6, 0.4, "Still contested", 15, True, ACCENT, SANS)
    multiline(
        s, 9.0, 2.8, 3.6, 3.0,
        [
            "A4 — defaults stay classical until distributional CI wins",
            "P1 — Magellan vs MLGO path",
            "P2 / S5 — multi-DSL skills & profiler APIs",
        ],
        13, INK, SANS, gap=12,
    )
    page_footer(s, n, total)


def s_commercial(prs, n, total):
    s = blank(prs)
    kicker(s, 0.7, 0.35, "13  —  Commercialization §5.7")
    textbox(s, 0.7, 0.75, 12, 0.5, "Ship the hybrid — solve P1–P23.", 28, True, INK, DISPLAY)
    bands = [
        (0.5, TEAL, "Architecture", "Contract · memory\ntopology · when-to-run\nmulti-DSL · FSM"),
        (4.7, ACCENT, "Trust & ops", "Oracles · ownership\nversioning · HITL\nDR · A/B gates"),
        (8.9, INK, "Business", "SKU · eval · pricing\ntenancy · IP\ntokens / capability"),
    ]
    for x, col, title, body in bands:
        round_rect(s, x, 1.55, 3.9, 3.2, CARD, LINE)
        rect(s, x, 1.55, 3.9, 0.12, col)
        textbox(s, x + 0.25, 1.9, 3.4, 0.45, title, 18, True, INK, DISPLAY)
        multiline(s, x + 0.25, 2.5, 3.4, 2.0, body.split("\n"), 15, MUTED, SANS, gap=8)
    textbox(
        s, 0.7, 5.05, 12, 0.4,
        "P23 verdict: tokens / latency / capability shape the SKU — falsify always-on LLM-as-opt, not frozen artifacts.",
        14, True, ACCENT, SANS,
    )
    # three-step path
    steps = [("Search", "Amdahl-hot · oracle-gated"), ("Freeze", "ACF / kernel → VCS"), ("Serve", "Zero LLM calls")]
    for i, (t, b) in enumerate(steps):
        x = 0.7 + i * 4.2
        round_rect(s, x, 5.6, 3.9, 1.15, SOFT_TEAL if i < 2 else SOFT_COPPER)
        textbox(s, x + 0.2, 5.7, 3.5, 0.35, f"{i+1}. {t}", 16, True, INK, DISPLAY)
        textbox(s, x + 0.2, 6.15, 3.5, 0.35, b, 13, False, MUTED, SANS)
    page_footer(s, n, total)


def s_org_questions(prs, n, total):
    s = blank(prs)
    kicker(s, 0.7, 0.4, "14  —  If you adopt this")
    textbox(s, 0.7, 0.85, 12, 0.55, "Questions before tools.", 28, True, INK, DISPLAY)
    qs = [
        ("01", "Online or offline first?", "Specialize workloads, or evolve the compiler once?"),
        ("02", "What is the oracle?", "Alive2, golden kernels, OpInfo, serving A/B — who owns false negatives?"),
        ("03", "Which agent contract IR?", "LLVM, MLIR, Triton, StableHLO, Tile — pick intentionally."),
        ("04", "Who owns agent artifacts?", "Named maintainer six months after merge."),
        ("05", "Token / latency budget?", "Freeze path vs always-on — what $/ %gain is allowed?"),
    ]
    for i, (num, q, a) in enumerate(qs):
        y = 1.55 + i * 1.0
        textbox(s, 0.7, y, 1.0, 0.4, num, 18, True, ACCENT, DISPLAY)
        textbox(s, 1.8, y, 10.5, 0.35, q, 16, True, INK, DISPLAY)
        textbox(s, 1.8, y + 0.38, 10.5, 0.35, a, 13, False, MUTED, SANS)
    page_footer(s, n, total)


def s_close(prs):
    s = blank(prs)
    rect(s, 0, 0, W, H, INK)
    rect(s, 0, 0, 0.18, H, ACCENT)
    textbox(s, 0.9, 1.6, 11, 0.4, "ONE-PAGE CHECK", 12, True, ACCENT, SANS)
    textbox(
        s, 0.9, 2.1, 11.5, 1.2,
        "Predict the architecture.\nName the four jobs.\nSeparate signal from noise.",
        28, True, WHITE, DISPLAY,
    )
    multiline(
        s, 0.9, 4.0, 11, 2.0,
        [
            "Architecture → hybrid control / data / codesign feedback",
            "Jobs → online · offline · oracle review · bring-up",
            "Signal → Tier A + CONFLICTS settlement watches",
            "Commercial → §5.7 P1–P23 · freeze artifacts · budget tokens",
        ],
        16, RGBColor(0xC8, 0xC3, 0xBA), SANS, gap=10,
    )
    textbox(
        s, 0.9, 6.4, 11, 0.4,
        "Rebuild:  python3 publish/build_pptx.py && python3 publish/build_visual.py",
        12, False, MUTED, SANS,
    )


def build() -> Path:
    OUT.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)

    slides = [
        s_title,
        s_thesis,
        s_not_replace,
        s_architecture,
        s_four_jobs,
        s_era,
        s_ideas_grid,
        s_stack,
        s_conflicts,
        s_codesign,
        s_roadmap,
        s_evidence_chart,
        s_claims,
        s_commercial,
        s_org_questions,
        s_close,
    ]
    total = len(slides)
    # title + close have custom chrome; numbered content slides get footers inside
    for i, fn in enumerate(slides, 1):
        if fn in (s_title, s_close):
            fn(prs)
        else:
            fn(prs, i, total)

    out = OUT / PPTX_NAME
    prs.save(str(out))
    return out


if __name__ == "__main__":
    path = build()
    print(f"Wrote {path} ({path.stat().st_size // 1024} KiB, editorial deck)")
