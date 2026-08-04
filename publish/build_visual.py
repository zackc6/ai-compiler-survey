#!/usr/bin/env python3
"""Render diagram-first survey visuals (PNG posters + visual PPTX) into publish/out.

Not table/text decks — each piece is one composition: architecture, orbit,
timeline path, ladder, or constellation.
"""

from __future__ import annotations

import math
from pathlib import Path

from PIL import Image, ImageDraw, ImageFilter, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = Path(__file__).resolve().parent / "out"
VIS = OUT / "visual"
PPTX_NAME = "next-gen-ai-compiler-survey-visual.pptx"

# Night-forge palette (avoid purple glow / cream+terracotta / broadsheet)
BG0 = (10, 14, 20)
BG1 = (18, 28, 36)
STEEL = (72, 140, 168)
AMBER = (232, 168, 56)
EMBER = (196, 96, 48)
MIST = (214, 220, 226)
MUTED = (120, 132, 144)
WHITE = (245, 247, 250)
INK = (8, 10, 14)

W, H = 1920, 1080  # 16:9 poster / slide

FONT_DIR_CANDIDATES = [
    Path("/usr/share/fonts/truetype/dejavu"),
    Path("/usr/share/fonts/truetype/croscore"),
    Path("/usr/share/fonts/truetype/macos"),
]


def find_font(names: list[str], size: int) -> ImageFont.FreeTypeFont:
    for d in FONT_DIR_CANDIDATES:
        for name in names:
            p = d / name
            if p.exists():
                return ImageFont.truetype(str(p), size)
    return ImageFont.load_default()


F_DISPLAY = lambda s: find_font(["Tinos-Bold.ttf", "DejaVuSerif-Bold.ttf"], s)
F_BODY = lambda s: find_font(["DejaVuSans.ttf", "Arimo-Regular.ttf"], s)
F_BOLD = lambda s: find_font(["DejaVuSans-Bold.ttf", "Arimo-Bold.ttf"], s)
F_MONO = lambda s: find_font(["JetBrainsMono-Regular.ttf", "DejaVuSansMono.ttf"], s)


def canvas() -> tuple[Image.Image, ImageDraw.ImageDraw]:
    """Solid night field + soft radial wash. No grid — grids hurt text readability."""
    img = Image.new("RGB", (W, H), BG0)
    overlay = Image.new("RGB", (W, H), BG0)
    od = ImageDraw.Draw(overlay)
    cx, cy = int(W * 0.55), int(H * 0.42)
    for r in range(1100, 0, -16):
        t = r / 1100
        c = (
            int(BG0[0] * t + BG1[0] * (1 - t)),
            int(BG0[1] * t + (BG1[1] + 6) * (1 - t)),
            int(BG0[2] * t + (BG1[2] + 10) * (1 - t)),
        )
        od.ellipse([cx - r, cy - r, cx + r, cy + r], fill=c)
    img = Image.blend(img, overlay, 0.75)
    draw = ImageDraw.Draw(img, "RGBA")
    return img, draw


def text(draw, xy, s, font, fill=MIST, anchor="lt"):
    draw.text(xy, s, font=font, fill=fill, anchor=anchor)


def brand_mark(draw, x=72, y=48):
    draw.rectangle([x, y, x + 10, y + 42], fill=AMBER)
    text(draw, (x + 28, y + 6), "NEXT-GEN AI COMPILER SURVEY", F_BOLD(22), AMBER)


def save(img: Image.Image, name: str) -> Path:
    VIS.mkdir(parents=True, exist_ok=True)
    path = VIS / name
    img.save(path, "PNG", optimize=True)
    print(f"wrote {path.relative_to(ROOT)} ({path.stat().st_size // 1024} KB)")
    return path


# ── compositions ────────────────────────────────────────────────────


def poster_title() -> Path:
    img, draw = canvas()
    brand_mark(draw)
    # dominant brand wordmark
    text(draw, (96, 280), "AGENTIC", F_DISPLAY(118), WHITE)
    text(draw, (96, 410), "COMPILER", F_DISPLAY(118), AMBER)
    # short supporting line only
    text(
        draw,
        (100, 580),
        "Agents search · Compilers decide · Silicon feeds back",
        F_BODY(32),
        MIST,
    )
    # visual anchor: interlocking rings (control / data)
    cx, cy = 1480, 540
    for i, (r, col, w) in enumerate(
        [(260, STEEL + (90,), 14), (190, AMBER + (120,), 10), (120, EMBER + (100,), 8)]
    ):
        bbox = [cx - r, cy - r, cx + r, cy + r]
        draw.ellipse(bbox, outline=col, width=w)
    text(draw, (cx, cy - 18), "control", F_BOLD(22), STEEL, "mm")
    text(draw, (cx, cy + 18), "data", F_BOLD(22), AMBER, "mm")
    text(draw, (96, 980), "Prediction horizon  ·  2027–28 and ~5 years", F_MONO(20), MUTED)
    return save(img, "01-title-agentic-compiler.png")


def poster_hybrid() -> Path:
    img, draw = canvas()
    brand_mark(draw)
    text(draw, (96, 120), "THE CONTRACT", F_BOLD(18), STEEL)
    text(draw, (96, 160), "Hybrid, not replacement", F_DISPLAY(64), WHITE)

    # two interlocking planes as big shapes
    left = [(120, 320), (880, 320), (820, 900), (120, 900)]
    right = [(1040, 320), (1800, 320), (1800, 900), (1100, 900)]
    draw.polygon(left, fill=STEEL + (45,))
    draw.polygon(right, fill=AMBER + (40,))
    draw.line(left + [left[0]], fill=STEEL, width=3)
    draw.line(right + [right[0]], fill=AMBER, width=3)

    text(draw, (180, 380), "AGENTS", F_DISPLAY(48), STEEL)
    for i, line in enumerate(
        ["semantic search", "orchestration", "artifact synthesis", "bring-up loops"]
    ):
        text(draw, (180, 480 + i * 70), line, F_BODY(30), MIST)

    text(draw, (1180, 380), "COMPILERS", F_DISPLAY(48), AMBER)
    for i, line in enumerate(
        ["lowering", "legality", "measurement", "admit / fallback"]
    ):
        text(draw, (1180, 480 + i * 70), line, F_BODY(30), MIST)

    # bridge arrow
    draw.polygon(
        [(860, 560), (1060, 560), (1060, 540), (1140, 600), (1060, 660), (1060, 640), (860, 640)],
        fill=WHITE + (180,),
    )
    return save(img, "02-hybrid-contract.png")


def poster_timeline() -> Path:
    img, draw = canvas()
    brand_mark(draw)
    text(draw, (96, 120), "ERA PATH", F_BOLD(18), STEEL)
    text(draw, (96, 160), "How we got here", F_DISPLAY(64), WHITE)

    eras = [
        (2018, "DL compilers", "TVM · Ansor · MLIR", STEEL),
        (2021, "MLGO / RL gyms", "inlining · regalloc", (96, 160, 140)),
        (2023, "LLM enters IR", "pass lists · Meta LLM", AMBER),
        (2025, "Agentic hybrid", "kernels · heuristics · verify", EMBER),
    ]
    y = 520
    xs = [220, 640, 1060, 1480]
    # path
    draw.line([(xs[0], y), (xs[-1], y)], fill=MUTED + (120,), width=6)
    for i, ((year, title, sub, col), x) in enumerate(zip(eras, xs)):
        draw.ellipse([x - 28, y - 28, x + 28, y + 28], fill=col)
        draw.ellipse([x - 12, y - 12, x + 12, y + 12], fill=BG0)
        text(draw, (x, y - 120), str(year), F_DISPLAY(40), WHITE, "mm")
        text(draw, (x, y + 80), title, F_BOLD(26), col, "mm")
        text(draw, (x, y + 130), sub, F_BODY(20), MUTED, "mm")
        if i < len(xs) - 1:
            # chevron
            nx = xs[i + 1]
            mid = (x + nx) // 2
            draw.polygon(
                [(mid - 18, y - 14), (mid + 18, y), (mid - 18, y + 14)],
                fill=MUTED,
            )
    return save(img, "03-era-path.png")


def poster_four_jobs() -> Path:
    img, draw = canvas()
    brand_mark(draw)
    text(draw, (96, 120), "AGENT WORKLOAD", F_BOLD(18), STEEL)
    text(draw, (96, 160), "Four jobs", F_DISPLAY(64), WHITE)

    cx, cy = 960, 620
    R = 320
    jobs = [
        ("a", "Online\nspecialize", "CompileIQ · GEAK\nACCLAIM", STEEL, -135),
        ("b", "Offline\nevolve", "Magellan\nAlphaEvolve", AMBER, -45),
        ("c", "Oracle\nreview", "Archer\nLLVM agents", (120, 180, 160), 45),
        ("d", "Bring-up /\ncodesign", "TritorX\nKernelEvolve", EMBER, 135),
    ]
    # core
    draw.ellipse([cx - 110, cy - 110, cx + 110, cy + 110], fill=BG1, outline=WHITE, width=3)
    text(draw, (cx, cy - 18), "agentic", F_BOLD(22), WHITE, "mm")
    text(draw, (cx, cy + 18), "compiler", F_BOLD(22), AMBER, "mm")

    for letter, title, who, col, ang in jobs:
        rad = math.radians(ang)
        x = cx + int(R * math.cos(rad))
        y = cy + int(R * math.sin(rad))
        draw.line([(cx, cy), (x, y)], fill=col + (140,), width=4)
        draw.ellipse([x - 130, y - 130, x + 130, y + 130], fill=col + (35,), outline=col, width=4)
        text(draw, (x, y - 55), letter, F_DISPLAY(42), col, "mm")
        # title lines
        for j, line in enumerate(title.split("\n")):
            text(draw, (x, y - 10 + j * 28), line, F_BOLD(22), WHITE, "mm")
        for j, line in enumerate(who.split("\n")):
            text(draw, (x, y + 55 + j * 22), line, F_BODY(16), MUTED, "mm")
    return save(img, "04-four-jobs-orbit.png")


def poster_stack() -> Path:
    img, draw = canvas()
    brand_mark(draw)
    text(draw, (96, 110), "STACK PRESSURE", F_BOLD(18), STEEL)
    text(draw, (96, 150), "Where agents push", F_DISPLAY(56), WHITE)

    layers = [
        ("1 Framework", "graphs · hot ops", STEEL),
        ("2 Kernel DSL", "Triton · Helion · Tile", (80, 150, 170)),
        ("3 Portable IR", "MLIR · StableHLO", (90, 130, 150)),
        ("4 Mid / back", "passes · heuristics · ACF", AMBER),
        ("5 Oracles", "tests · Alive2 · profilers", EMBER),
        ("6 Artifacts", "ACF · kernels · memory", (180, 120, 70)),
        ("7 Serving", "specialize hot paths", (140, 100, 90)),
        ("8 Silicon / sim", "bring-up feedback", (160, 80, 60)),
    ]
    x0, y0, lw, lh, gap = 180, 260, 980, 78, 10
    for i, (name, sub, col) in enumerate(layers):
        y = y0 + i * (lh + gap)
        draw.rounded_rectangle([x0, y, x0 + lw, y + lh], radius=12, fill=col + (50,), outline=col, width=2)
        text(draw, (x0 + 28, y + 18), name, F_BOLD(26), WHITE)
        text(draw, (x0 + 420, y + 22), sub, F_BODY(22), MIST)

    # agent arrow column
    ax = 1320
    draw.rounded_rectangle([ax, 260, ax + 480, 960], radius=18, fill=BG1, outline=AMBER, width=3)
    text(draw, (ax + 40, 300), "Agent jobs", F_BOLD(28), AMBER)
    arrows = [
        ("(a) Online", "1–2–4–5–7"),
        ("(b) Offline", "4 → artifacts"),
        ("(c) Review", "4–5–6"),
        ("(d) Codesign", "2–5–8"),
    ]
    for i, (a, b) in enumerate(arrows):
        y = 400 + i * 120
        text(draw, (ax + 40, y), a, F_BOLD(26), WHITE)
        text(draw, (ax + 40, y + 40), b, F_MONO(22), MUTED)
    return save(img, "05-stack-pressure.png")


def poster_codesign() -> Path:
    img, draw = canvas()
    brand_mark(draw)
    text(draw, (96, 120), "HW–SW LOOP", F_BOLD(18), STEEL)
    text(draw, (96, 160), "Coverage, then peak", F_DISPLAY(56), WHITE)

    # ascending steps
    steps = [
        ("SIM", "future device\nQEMU / draft ISA", STEEL, 220),
        ("COVERAGE", "TritorX\nops that run", (90, 160, 150), 420),
        ("PERF", "KernelEvolve\nhetero search", AMBER, 620),
        ("FEEDBACK", "ISA / dialect\nRFC to humans", EMBER, 820),
    ]
    base_y = 880
    for i, (title, body, col, h) in enumerate(steps):
        x = 180 + i * 420
        top = base_y - h
        draw.polygon(
            [(x, base_y), (x + 340, base_y), (x + 340, top + 40), (x + 300, top), (x, top)],
            fill=col + (55,),
            outline=col,
        )
        draw.line([(x, base_y), (x + 340, base_y), (x + 340, top + 40), (x + 300, top), (x, top), (x, base_y)], fill=col, width=3)
        text(draw, (x + 30, top + 50), title, F_DISPLAY(36), WHITE)
        for j, line in enumerate(body.split("\n")):
            text(draw, (x + 30, top + 120 + j * 36), line, F_BODY(24), MIST)
        if i < len(steps) - 1:
            draw.polygon(
                [(x + 350, base_y - h // 2 - 20), (x + 400, base_y - h // 2), (x + 350, base_y - h // 2 + 20)],
                fill=MUTED,
            )
    text(draw, (96, 980), "Not autonomous EDA — agents stress compilers; humans own tape-out", F_BODY(22), MUTED)
    return save(img, "06-codesign-ladder.png")


def poster_conflicts() -> Path:
    img, draw = canvas()
    brand_mark(draw)
    text(draw, (96, 120), "PRODUCTIVE TENSION", F_BOLD(18), STEEL)
    text(draw, (96, 160), "Do not average these", F_DISPLAY(56), WHITE)

    pairs = [
        ("C1", "Magellan\nC++ heuristics", "MLGO\nneural advisors"),
        ("C2", "Vendor\nspeedups", "KernelBench-X\nceilings"),
        ("C3", "Free rewrite", "Advisory\n+ admit"),
        ("C9", "Coverage\nfirst", "Peak perf\nfirst"),
        ("C10", "Compiler\ncodesign", "Autonomous\nchip design"),
    ]
    for i, (cid, left, right) in enumerate(pairs):
        y = 280 + i * 140
        # left
        draw.rounded_rectangle([140, y, 780, y + 110], radius=14, fill=STEEL + (40,), outline=STEEL, width=2)
        text(draw, (160, y + 20), cid, F_BOLD(22), AMBER)
        for j, line in enumerate(left.split("\n")):
            text(draw, (260, y + 20 + j * 36), line, F_BODY(26), WHITE)
        # vs node
        draw.ellipse([900, y + 20, 1020, y + 90], fill=BG1, outline=AMBER, width=3)
        text(draw, (960, y + 55), "vs", F_BOLD(24), AMBER, "mm")
        # right
        draw.rounded_rectangle([1140, y, 1780, y + 110], radius=14, fill=EMBER + (40,), outline=EMBER, width=2)
        for j, line in enumerate(right.split("\n")):
            text(draw, (1180, y + 20 + j * 36), line, F_BODY(26), WHITE)
    return save(img, "07-conflicts-tension.png")


def poster_constellation() -> Path:
    img, draw = canvas()
    brand_mark(draw)
    text(draw, (96, 120), "VISION MAP", F_BOLD(18), STEEL)
    text(draw, (96, 160), "Public agendas in orbit", F_DISPLAY(52), WHITE)

    # center survey
    cx, cy = 960, 580
    draw.ellipse([cx - 160, cy - 160, cx + 160, cy + 160], fill=AMBER + (50,), outline=AMBER, width=4)
    text(draw, (cx, cy - 24), "this", F_BOLD(28), WHITE, "mm")
    text(draw, (cx, cy + 16), "survey", F_BOLD(28), WHITE, "mm")

    nodes = [
        ("Compiler 2.0\n/ MOCHA", -110, 280, STEEL),
        ("New Compiler\nStack", -20, 280, (90, 160, 170)),
        ("Compiler.next", 70, 280, AMBER),
        ("Kernel agents\nKernelEvolve", 160, 300, EMBER),
        ("Magellan /\nAlphaEvolve", 200, 220, (180, 140, 80)),
        ("ACCLAIM /\nhybrid loops", -160, 240, (100, 170, 140)),
    ]
    for label, ang, dist, col in nodes:
        rad = math.radians(ang)
        x = cx + int(dist * math.cos(rad))
        y = cy + int(dist * math.sin(rad))
        draw.line([(cx, cy), (x, y)], fill=col + (100,), width=3)
        draw.ellipse([x - 110, y - 70, x + 110, y + 70], fill=BG1, outline=col, width=3)
        for j, line in enumerate(label.split("\n")):
            text(draw, (x, y - 18 + j * 28), line, F_BOLD(20), WHITE, "mm")
    return save(img, "08-vision-constellation.png")


def poster_architecture() -> Path:
    """§5.1 hybrid stack: intent → control plane + substrate → data plane → sinks + codesign."""
    img, draw = canvas()
    brand_mark(draw)
    text(draw, (96, 88), "§5.1 ARCHITECTURE", F_BOLD(18), STEEL)
    text(draw, (96, 120), "Hybrid agentic compiler", F_DISPLAY(48), WHITE)
    text(
        draw,
        (96, 178),
        "Agents orchestrate · classical compilers execute · silicon feeds the next dialect/ISA RFC",
        F_BODY(20),
        MUTED,
    )

    # Intent pill
    draw.rounded_rectangle([560, 220, 1360, 268], radius=22, fill=BG1, outline=MIST + (80,), width=2)
    text(draw, (960, 244), "HUMAN / PRODUCT INTENT  ·  NL at edge · policy · budget", F_BOLD(18), MIST, "mm")

    # Control plane
    draw.rounded_rectangle([80, 290, 1840, 560], radius=16, fill=STEEL + (36,), outline=STEEL, width=3)
    text(draw, (110, 308), "AGENT CONTROL PLANE", F_BOLD(20), STEEL)
    text(draw, (480, 312), "lab → CI-gated → hot-path specialize", F_BODY(18), MUTED)

    jobs = [
        ("(a)", "ONLINE", "propose · measure · admit"),
        ("(b)", "OFFLINE", "heuristics → C++ / MLGO"),
        ("(c)", "ENG.", "oracle PR / src edit"),
        ("(d)", "BRING-UP", "coverage → perf · sim+Si"),
    ]
    for i, (letter, title, sub) in enumerate(jobs):
        x = 110 + i * 425
        draw.rounded_rectangle([x, 350, x + 400, 455], radius=12, fill=BG1, outline=AMBER, width=2)
        text(draw, (x + 20, 368), letter, F_BOLD(22), AMBER)
        text(draw, (x + 80, 366), title, F_BOLD(24), WHITE)
        text(draw, (x + 20, 412), sub, F_BODY(18), MIST)

    draw.rounded_rectangle([110, 475, 1810, 535], radius=10, fill=AMBER + (32,), outline=AMBER, width=2)
    text(draw, (130, 488), "SUBSTRATE", F_BOLD(18), AMBER)
    text(
        draw,
        (280, 488),
        "workflow compile (FlowCompile)  ·  ADG analysis (AgentFlow)  ·  freeze/amortize (Auto)  ·  hetero place",
        F_BODY(18),
        WHITE,
    )

    # Contract arrow band
    draw.polygon([(900, 575), (1020, 575), (1020, 598), (960, 640), (900, 598)], fill=WHITE + (150,))
    text(draw, (1040, 590), "typed tools · admit records · oracles · FSM / plan bounds", F_MONO(17), MUTED)

    # Data plane
    draw.rounded_rectangle([80, 660, 1840, 800], radius=16, fill=AMBER + (28,), outline=AMBER, width=3)
    text(draw, (110, 678), "CLASSICAL DATA PLANE", F_BOLD(20), AMBER)
    text(draw, (520, 682), "default path stays until CI proves agents", F_BODY(18), MUTED)
    text(draw, (110, 720), "Framework → Inductor · XLA · MLIR · Triton · Helion · Tile · CuTe · device libs", F_BODY(22), WHITE)
    text(draw, (110, 760), "legality · lowering · cost models · golden / Alive2 / OpInfo · admit / fallback", F_BODY(18), MIST)

    # Three sinks + codesign feedback
    sinks = [
        ("GPU / NPU / ASIC", "sim → silicon", STEEL),
        ("VCS artifacts", "ACF · kernels · heuristics · memory", (180, 120, 70)),
        ("Serving runtime", "specialize hot paths · freeze replay", EMBER),
    ]
    for i, (title, sub, col) in enumerate(sinks):
        x = 80 + i * 420
        draw.rounded_rectangle([x, 830, x + 400, 920], radius=12, fill=col + (38,), outline=col, width=2)
        text(draw, (x + 22, 848), title, F_BOLD(20), WHITE)
        text(draw, (x + 22, 882), sub, F_BODY(16), MIST)

    draw.rounded_rectangle([1360, 830, 1840, 920], radius=12, fill=EMBER + (28,), outline=EMBER, width=2)
    text(draw, (1380, 848), "HW–SW CODESIGN FEEDBACK", F_BOLD(18), EMBER)
    text(draw, (1380, 882), "failures → ISA / dialect RFCs  ·  C10", F_BODY(16), MIST)

    text(
        draw,
        (80, 980),
        "Invariant: LLM outputs guide search — they must not silently define unchecked executable behavior",
        F_BODY(18),
        MUTED,
    )
    return save(img, "02b-architecture-51.png")


def poster_horizon() -> Path:
    """Roadmap architecture evolution: Today → Horizon A → Horizon B as stacked planes."""
    img, draw = canvas()
    brand_mark(draw)
    text(draw, (96, 80), "§5.5 ROADMAP", F_BOLD(18), STEEL)
    text(draw, (96, 112), "Architecture evolution", F_DISPLAY(48), WHITE)
    text(
        draw,
        (96, 168),
        "Left → right: agents stay on the control plane; the data plane never goes away",
        F_BODY(20),
        MUTED,
    )

    # Timeline spine
    draw.line([(220, 250), (1700, 250)], fill=MUTED + (110,), width=4)
    stages = [
        (
            "TODAY",
            "2025–26",
            STEEL,
            "CONTROL",
            ["Ad-hoc agent loops", "chat / tools", "GEAK · ACCLAIM · Magellan"],
            "DATA",
            ["MLIR / Triton / Inductor", "still the default", "GPU-mostly"],
            "CODESIGN",
            ["mostly absent", "manual bring-up"],
        ),
        (
            "HORIZON A",
            "2027–28",
            AMBER,
            "CONTROL",
            ["Jobs (a–d) productized", "CI-gated specialize", "typed tools + oracles"],
            "DATA",
            ["multi-DSL + fingerprints", "tool APIs exposed", "admit / fallback"],
            "CODESIGN",
            ["early: sim + 1st Si", "→ ISA / dialect RFCs"],
        ),
        (
            "HORIZON B",
            "~2029–31",
            EMBER,
            "CONTROL",
            ["control plane compiled", "ADG · freeze · place", "workflow amortize"],
            "DATA",
            ["multi-backend fleets", "ACF / heuristics / memory", "as VCS artifacts"],
            "CODESIGN",
            ["steady pre-Si loop", "not autonomous EDA"],
        ),
    ]

    for i, (label, years, col, c_lbl, c_lines, d_lbl, d_lines, h_lbl, h_lines) in enumerate(stages):
        x = 90 + i * 610
        # node on spine
        draw.ellipse([x + 250, 230, x + 290, 270], fill=col, outline=WHITE, width=2)
        if i < 2:
            draw.polygon([(x + 545, 240), (x + 575, 250), (x + 545, 260)], fill=MUTED)

        # header
        text(draw, (x + 36, 290), label, F_BOLD(20), col)
        text(draw, (x + 36, 322), years, F_DISPLAY(34), WHITE)

        # three stacked bands
        bands = [
            (370, c_lbl, c_lines, STEEL),
            (560, d_lbl, d_lines, AMBER),
            (750, h_lbl, h_lines, EMBER),
        ]
        for y0, blabel, lines, bcol in bands:
            draw.rounded_rectangle([x, y0, x + 560, y0 + 170], radius=14, fill=bcol + (30,), outline=bcol, width=2)
            text(draw, (x + 24, y0 + 16), blabel, F_BOLD(16), bcol)
            for j, line in enumerate(lines):
                text(draw, (x + 24, y0 + 52 + j * 34), line, F_BODY(20), WHITE)

    text(
        draw,
        (96, 990),
        "Artifact store grows → binaries → +ACF/hints → +evolved C++ / verified kernels / bring-up corpora → +frozen agent workflows",
        F_BODY(18),
        MUTED,
    )
    return save(img, "09-architecture-evolution.png")


def poster_falsifiers() -> Path:
    img, draw = canvas()
    brand_mark(draw)
    text(draw, (96, 120), "KILL CRITERIA", F_BOLD(18), EMBER)
    text(draw, (96, 160), "What would falsify this", F_DISPLAY(52), WHITE)

    items = [
        ("01", "Default AI stack ships with no classical admit / fallback"),
        ("02", "Both Magellan-class synthesis and MLGO advisors vanish"),
        ("03", "Kernel agents forever lose to eager on fusion-heavy suites"),
    ]
    for i, (num, line) in enumerate(items):
        y = 340 + i * 180
        draw.rounded_rectangle([140, y, 1780, y + 140], radius=20, fill=BG1, outline=EMBER, width=3)
        text(draw, (200, y + 45), num, F_DISPLAY(48), EMBER)
        text(draw, (360, y + 50), line, F_BODY(32), WHITE)
    return save(img, "10-falsifiers.png")


def poster_commercial() -> Path:
    """§5.7 commercialization — three problem bands as one composition."""
    img, draw = canvas()
    brand_mark(draw)
    text(draw, (96, 120), "COMMERCIAL PRACTICE", F_BOLD(18), STEEL)
    text(draw, (96, 160), "P1–P23 problem bands", F_DISPLAY(52), WHITE)

    bands = [
        (
            280,
            STEEL,
            "ARCHITECTURE",
            ["P1 contract", "P2 memory", "P3 topology", "P4 when-to-run", "P7 multi-DSL", "P22 FSM bounds"],
        ),
        (
            520,
            AMBER,
            "TRUST & OPS",
            ["P5 oracles", "P6 ownership", "P14 versioning", "P16 HITL", "P19 DR", "P20 A/B"],
        ),
        (
            760,
            EMBER,
            "BUSINESS",
            ["P8 SKU", "P9 eval", "P10 pricing", "P11 tenancy", "P13 IP", "P23 tokens/capability"],
        ),
    ]
    for y, col, title, items in bands:
        draw.rounded_rectangle([120, y, 1800, y + 200], radius=18, fill=col + (28,), outline=col, width=3)
        text(draw, (160, y + 30), title, F_BOLD(26), col)
        for i, item in enumerate(items):
            x = 160 + (i % 6) * 270
            yy = y + 90 + (i // 6) * 50
            draw.ellipse([x, yy + 8, x + 14, yy + 22], fill=col)
            text(draw, (x + 28, yy), item, F_BODY(22), WHITE)
    text(
        draw,
        (120, 1000),
        "Lean: typed admit traces · freeze artifacts · customer owns outputs · budget tokens per %gain",
        F_BODY(22),
        MUTED,
    )
    return save(img, "11-commercial-bands.png")


def poster_resource_envelope() -> Path:
    """P23 — tokens / inference / capability as three gauges + verdict."""
    img, draw = canvas()
    brand_mark(draw)
    text(draw, (96, 120), "RESOURCE ENVELOPE", F_BOLD(18), STEEL)
    text(draw, (96, 160), "Tokens · latency · capability", F_DISPLAY(48), WHITE)

    gauges = [
        (320, "TOKENS", "Always-on online\nburns OpEx", "Freeze + Amdahl\nbudgets", EMBER),
        (960, "INFERENCE", "Hours-scale search\n≠ chat TTFT", "CI / overnight\nlab vs product", AMBER),
        (1600, "CAPABILITY", "Weak one-shot\nunsafe free IR\nfragile tools", "Search + oracles\nclassical data plane", STEEL),
    ]
    for cx, title, problem, fix, col in gauges:
        draw.ellipse([cx - 200, 340, cx + 200, 740], outline=col, width=8)
        draw.ellipse([cx - 150, 390, cx + 150, 690], fill=col + (35,))
        text(draw, (cx, 430), title, F_BOLD(28), col, "mm")
        for j, line in enumerate(problem.split("\n")):
            text(draw, (cx, 500 + j * 32), line, F_BODY(20), MIST, "mm")
        text(draw, (cx, 640), "→", F_BOLD(28), WHITE, "mm")
        for j, line in enumerate(fix.split("\n")):
            text(draw, (cx, 680 + j * 28), line, F_BOLD(18), WHITE, "mm")

    # verdict bar
    draw.rounded_rectangle([120, 860, 1800, 1000], radius=16, fill=BG1, outline=AMBER, width=3)
    text(draw, (160, 890), "VERDICT", F_BOLD(22), AMBER)
    text(
        draw,
        (160, 935),
        "Shapes the SKU — does not kill hybrid. Falsifies LLM-as-opt every build. Not agentic control planes.",
        F_BODY(26),
        WHITE,
    )
    return save(img, "12-resource-envelope.png")


def poster_freeze_amortize() -> Path:
    """Commercial path: spend tokens in CI → ship artifact → zero LLM at serve."""
    img, draw = canvas()
    brand_mark(draw)
    text(draw, (96, 120), "COMMERCIAL PATH", F_BOLD(18), STEEL)
    text(draw, (96, 160), "Amortize the LLM away", F_DISPLAY(52), WHITE)

    stages = [
        (200, "1 SEARCH", "Frontier tokens\nAmdahl-hot only\noracle-gated", STEEL),
        (700, "2 FREEZE", "ACF / kernel /\nheuristic → VCS\nadmit record", AMBER),
        (1200, "3 SERVE", "Classical path\nmillions of infer\nzero LLM calls", EMBER),
    ]
    for x, title, body, col in stages:
        draw.rounded_rectangle([x, 360, x + 420, 780], radius=20, fill=col + (40,), outline=col, width=4)
        text(draw, (x + 40, 400), title, F_DISPLAY(36), WHITE)
        for j, line in enumerate(body.split("\n")):
            text(draw, (x + 40, 500 + j * 50), line, F_BODY(26), MIST)
    # arrows
    for x in (620, 1120):
        draw.polygon(
            [(x, 540), (x + 60, 570), (x, 600)],
            fill=MUTED,
        )
    text(
        draw,
        (200, 900),
        "Token & capability risk lives in eng CI — not in the customer critical path",
        F_BODY(26),
        MUTED,
    )
    return save(img, "13-freeze-amortize.png")


def build_pptx(image_paths: list[Path], name: str = PPTX_NAME) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for path in image_paths:
        s = prs.slides.add_slide(blank)
        # full-bleed image
        s.shapes.add_picture(str(path), Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    out = OUT / name
    prs.save(out)
    print(f"wrote {out.relative_to(ROOT)} ({out.stat().st_size // 1024} KB)")
    return out


def main() -> int:
    OUT.mkdir(parents=True, exist_ok=True)
    VIS.mkdir(parents=True, exist_ok=True)
    arch = poster_architecture()
    evo = poster_horizon()
    paths = [
        poster_title(),
        poster_hybrid(),
        arch,
        poster_timeline(),
        poster_four_jobs(),
        poster_stack(),
        poster_codesign(),
        poster_conflicts(),
        poster_constellation(),
        evo,
        poster_falsifiers(),
        poster_commercial(),
        poster_resource_envelope(),
        poster_freeze_amortize(),
    ]
    build_pptx(paths)
    # Focused 2-slide share deck for §5.1 + evolution
    build_pptx(
        [arch, evo],
        "architecture-51-and-evolution.pptx",
    )
    (VIS / "README.md").write_text(
        "# Survey visuals\n\nDiagram-first posters for the living survey "
        "(not table/text slides). Regenerated by `python3 publish/build_visual.py`.\n\n"
        "**Share these first:**\n"
        "- `02b-architecture-51.png` — SURVEY §5.1 hybrid stack\n"
        "- `09-architecture-evolution.png` — Today → Horizon A → Horizon B\n"
        "- Parent deck: `../next-gen-ai-compiler-survey-visual.pptx`\n"
        "- Focused 2-slide deck: `../architecture-51-and-evolution.pptx`\n\n"
        "Also includes commercialization (§5.7) and resource envelope (P23) posters.\n",
        encoding="utf-8",
    )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
